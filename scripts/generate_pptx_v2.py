#!/usr/bin/env python3
"""
PETITION ENGINE — PPTX Generator V2 (Premium OOXML)
Design DNA: SOI Premium Deck + Dabus Framework + Cerbasi Dossier
Fonts: Palatino Linotype (headings) + Garamond (body)
Palette: Navy #1B2A4A | Gold #B8860B | Cream #E8D5B7 | Body #3A3A3A

Uses raw OOXML XML injection for pixel-perfect positioning.
"""

import json
import os
import sys
import argparse
import html
import glob as globmod
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime

# ============================================================
# DESIGN DNA — Exact coordinates from QA analysis
# All values in POINTS (1pt = 12700 EMU)
# Slide: 720 × 405 pt (widescreen 16:9)
# ============================================================

# Conversion
def pt2emu(pts):
    return int(pts * 12700)

# Colors (hex strings for OOXML)
NAVY     = "1B2A4A"
GOLD     = "B8860B"
CREAM    = "E8D5B7"
WHITE    = "FFFFFF"
BODY_CLR = "3A3A3A"
CARD_BG  = "F5F0EB"
ALT_ROW  = "F7F5F2"
LT_GRAY  = "E8E8E8"
DARK_BG  = "0F172A"

# python-pptx color objects
C_NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
C_GOLD  = RGBColor(0xB8, 0x86, 0x0B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Typography
FONT_HEADING = "Palatino Linotype"
FONT_BODY    = "Garamond"

# Grid (pts)
SLIDE_W  = 720
SLIDE_H  = 405
MARGIN_L = 22
MARGIN_R = 22
CONTENT_W = 677       # 720 - 22 - 22 (actually 676, using 677)
TITLE_BAR_H = 58
TITLE_X  = 22
TITLE_Y  = 11
TITLE_W  = 677
TITLE_H  = 36
FOOTER_Y = 367
FOOTER_H = 29
BODY_START_Y = 72     # below title bar
BODY_END_Y = 360      # above footer

# Card layouts
CARD_2COL_W = 331
CARD_2COL_GAP = 15
CARD_3COL_W = 216
CARD_3COL_GAP = 14
CARD_PADDING = 10

# Slide dimensions in EMU
SLIDE_W_EMU = pt2emu(720)
SLIDE_H_EMU = pt2emu(405)

# Shape ID counter
_shape_id = [100]
def next_id():
    _shape_id[0] += 1
    return _shape_id[0]

def esc(text):
    """Escape text for XML."""
    return html.escape(str(text))


# ============================================================
# OOXML INJECTION HELPERS
# ============================================================
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_spTree(slide):
    """Get the shape tree element from a slide."""
    return slide._element.find(f'{{{NS_P}}}cSld/{{{NS_P}}}spTree')

def add_raw_shape(slide, xml_str):
    """Inject raw OOXML shape XML into slide."""
    spTree = get_spTree(slide)
    el = etree.fromstring(xml_str.encode('utf-8'))
    spTree.append(el)

def add_rect(slide, x, y, w, h, fill, name="Rect"):
    """Add a filled rectangle via OOXML."""
    sid = next_id()
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="{name}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
        <a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
    </p:sp>''')

def add_round_rect(slide, x, y, w, h, fill, border_color=None, name="Card"):
    """Add a rounded rectangle card."""
    sid = next_id()
    border = f'<a:ln w="6350"><a:solidFill><a:srgbClr val="{border_color}"><a:alpha val="30000"/></a:srgbClr></a:solidFill></a:ln>' if border_color else '<a:ln w="0"><a:noFill/></a:ln>'
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="{name}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 5000"/></a:avLst></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
        {border}
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
    </p:sp>''')

def add_circle(slide, x, y, d, fill, text="", text_color=WHITE, font_size=14, font=None):
    """Add a circle with centered text."""
    sid = next_id()
    f = font or FONT_HEADING
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Circle"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(d)}" cy="{pt2emu(d)}"/></a:xfrm>
        <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
        <a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" anchor="ctr" anchorCtr="1"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="{font_size*100}" b="1" dirty="0">
          <a:solidFill><a:srgbClr val="{text_color}"/></a:solidFill>
          <a:latin typeface="{f}"/></a:rPr>
          <a:t>{esc(text)}</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>''')

def add_textbox(slide, x, y, w, h, text, font_size=11, color=BODY_CLR,
                font=None, bold=False, italic=False, align="l",
                line_spacing=None, spc=None):
    """Add a text box with precise formatting."""
    sid = next_id()
    f = font or FONT_BODY
    b = ' b="1"' if bold else ""
    i = ' i="1"' if italic else ""
    sz = int(font_size * 100)
    lnSpc = f'<a:lnSpc><a:spcPts val="{line_spacing}"/></a:lnSpc>' if line_spacing else ""
    spacing = f' spc="{spc}"' if spc else ""
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="TB"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="t"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="{align}">{lnSpc}</a:pPr>
          <a:r><a:rPr lang="en-US" sz="{sz}"{b}{i}{spacing} dirty="0">
            <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
            <a:latin typeface="{f}" pitchFamily="34" charset="0"/></a:rPr>
            <a:t>{esc(text)}</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>''')

def add_multi_text(slide, x, y, w, h, paragraphs_data):
    """Add textbox with multiple paragraphs. paragraphs_data: list of dicts with text, size, color, bold, font, align, lnSpc."""
    sid = next_id()
    paras_xml = ""
    for p in paragraphs_data:
        f = p.get('font', FONT_BODY)
        sz = int(p.get('size', 11) * 100)
        b = ' b="1"' if p.get('bold') else ""
        i = ' i="1"' if p.get('italic') else ""
        color = p.get('color', BODY_CLR)
        align = p.get('align', 'l')
        lnSpc = f'<a:lnSpc><a:spcPts val="{p["lnSpc"]}"/></a:lnSpc>' if p.get('lnSpc') else ""
        spcAft = f'<a:spcAft><a:spcPts val="{p["spcAft"]}"/></a:spcAft>' if p.get('spcAft') else ""
        # Handle bullets
        buChar = f'<a:buChar char="{p["bullet"]}"/>' if p.get('bullet') else '<a:buNone/>'
        marL = f' marL="{p["marL"]}"' if p.get('marL') else ""
        indent = f' indent="{p["indent"]}"' if p.get('indent') else ""

        paras_xml += f'''<a:p><a:pPr algn="{align}"{marL}{indent}>{lnSpc}{spcAft}{buChar}</a:pPr>
          <a:r><a:rPr lang="en-US" sz="{sz}"{b}{i} dirty="0">
            <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
            <a:latin typeface="{f}" pitchFamily="34" charset="0"/></a:rPr>
            <a:t>{esc(p['text'])}</a:t></a:r></a:p>'''

    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="MTB"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="t"/>
        <a:lstStyle/>
        {paras_xml}
      </p:txBody>
    </p:sp>''')

def add_arrow(slide, x, y, w, h, fill=GOLD):
    """Add a right arrow shape."""
    sid = next_id()
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Arrow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="rightArrow"><a:avLst>
          <a:gd name="adj1" fmla="val 50000"/><a:gd name="adj2" fmla="val 50000"/>
        </a:avLst></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
        <a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
    </p:sp>''')

def add_gold_rule(slide, x, y, w, h=2):
    """Add a thin gold decorative line."""
    add_rect(slide, x, y, w, h, GOLD, "GoldRule")

def add_line(slide, x, y, w, h=1, color=GOLD, alpha=None):
    """Add a thin line (rect)."""
    sid = next_id()
    alpha_xml = f'<a:alpha val="{alpha}"/>' if alpha else ""
    add_raw_shape(slide, f'''
    <p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Line"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{pt2emu(x)}" y="{pt2emu(y)}"/>
                 <a:ext cx="{pt2emu(w)}" cy="{pt2emu(h)}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{color}">{alpha_xml}</a:srgbClr></a:solidFill>
        <a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
    </p:sp>''')


# ============================================================
# STANDARD SLIDE COMPONENTS
# ============================================================
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(int(color[:2],16), int(color[2:4],16), int(color[4:],16))

def add_title_bar(slide, title_text):
    """Navy title bar at top with gold Palatino title."""
    add_rect(slide, 0, 0, SLIDE_W, TITLE_BAR_H, NAVY, "TitleBar")
    add_textbox(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, title_text,
                font_size=28, color=GOLD, font=FONT_HEADING, bold=True)

def add_footer(slide, client_name, doc_label):
    """Footer at y=367, Garamond 9pt cream, centered."""
    add_textbox(slide, MARGIN_L, FOOTER_Y, CONTENT_W, FOOTER_H,
                f"{client_name} | {doc_label}",
                font_size=9, color=CREAM, font=FONT_BODY, align="ctr")


# ============================================================
# SLIDE BUILDERS
# ============================================================
def build_cover(prs, client_name, title, subtitle, doc_label, visa_type=""):
    """Cover slide with gold bars, centered layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Gold bar top
    add_rect(slide, 0, 0, SLIDE_W, 11, GOLD)

    # Title — uppercase, letter-spaced
    add_textbox(slide, MARGIN_L, 80, CONTENT_W, 60, title.upper(),
                font_size=36, color=GOLD, font=FONT_HEADING, bold=True,
                align="ctr", spc="400")

    # Subtitle
    if subtitle:
        add_textbox(slide, MARGIN_L, 150, CONTENT_W, 30, subtitle,
                    font_size=18, color=WHITE, font=FONT_HEADING, align="ctr")

    # Gold decorative line centered
    add_gold_rule(slide, (SLIDE_W - 288) / 2, 195, 288)

    # Client name
    add_textbox(slide, MARGIN_L, 210, CONTENT_W, 40, client_name,
                font_size=32, color=GOLD, font=FONT_HEADING, bold=True, align="ctr")

    # Visa type
    if visa_type:
        add_textbox(slide, MARGIN_L, 258, CONTENT_W, 25, visa_type,
                    font_size=14, color=WHITE, font=FONT_BODY, align="ctr")

    # Date
    add_textbox(slide, MARGIN_L, 290, CONTENT_W, 20,
                datetime.now().strftime("%B %Y"),
                font_size=12, color=CREAM, font=FONT_BODY, align="ctr")

    # Gold bar bottom
    add_rect(slide, 0, SLIDE_H - 11, SLIDE_W, 11, GOLD)

    add_footer(slide, client_name, doc_label)
    return slide


def build_toc(prs, sections, client_name, doc_label):
    """TOC with roman numerals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_title_bar(slide, "TABLE OF CONTENTS")

    romans = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X"]
    y = BODY_START_Y + 10
    for i, sec in enumerate(sections[:10]):
        rom = romans[i] if i < len(romans) else str(i+1)
        # Roman numeral
        add_textbox(slide, MARGIN_L, y, 40, 22, rom,
                    font_size=16, color=GOLD, font=FONT_HEADING, bold=True, align="ctr")
        # Section title
        add_textbox(slide, MARGIN_L + 50, y, CONTENT_W - 60, 22,
                    sec.get("title", ""),
                    font_size=12, color=WHITE, font=FONT_BODY, bold=True)
        # Subtitle
        if sec.get("subtitle"):
            add_textbox(slide, MARGIN_L + 50, y + 20, CONTENT_W - 60, 16,
                        sec["subtitle"],
                        font_size=9, color=CREAM, font=FONT_BODY, italic=True)
        y += 28

    add_footer(slide, client_name, doc_label)
    return slide


def build_section_divider(prs, title, subtitle, client_name, doc_label, number=""):
    """Section divider — navy bg, gold rules, centered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Gold rule top
    add_gold_rule(slide, (SLIDE_W - 600) / 2, 145, 600)

    # Volume/Section title
    display_title = f"{number}. {title}" if number else title
    add_textbox(slide, MARGIN_L, 155, CONTENT_W, 50, display_title.upper(),
                font_size=28, color=GOLD, font=FONT_HEADING, bold=True, align="ctr")

    # Subtitle
    if subtitle:
        add_textbox(slide, MARGIN_L, 215, CONTENT_W, 30, subtitle,
                    font_size=18, color=WHITE, font=FONT_BODY, align="ctr")

    # Gold rule bottom
    add_gold_rule(slide, (SLIDE_W - 600) / 2, 255, 600)

    add_footer(slide, client_name, doc_label)
    return slide


def build_content(prs, title, paragraphs, client_name, doc_label, bullets=None):
    """Content slide with title bar + body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    # Body paragraphs
    paras_data = []
    for p in (paragraphs or [])[:4]:
        paras_data.append({
            'text': p.strip(), 'size': 11, 'color': BODY_CLR,
            'font': FONT_BODY, 'lnSpc': 2200, 'spcAft': 600,
        })

    # Bullets
    for b in (bullets or []):
        paras_data.append({
            'text': b.strip(), 'size': 11, 'color': BODY_CLR,
            'font': FONT_BODY, 'lnSpc': 2200, 'spcAft': 400,
            'bullet': '\u2022', 'marL': '342900', 'indent': '-342900',
        })

    if paras_data:
        add_multi_text(slide, MARGIN_L + 7, BODY_START_Y + 8, CONTENT_W - 14,
                       BODY_END_Y - BODY_START_Y - 10, paras_data)

    add_footer(slide, client_name, doc_label)
    return slide


def build_metrics(prs, title, metrics, client_name, doc_label):
    """Stat cards — navy cards with gold numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    n = min(len(metrics), 4)
    if n <= 3:
        card_w, gap = CARD_3COL_W, CARD_3COL_GAP
    else:
        total_gap = (n - 1) * 14
        card_w = (CONTENT_W - 14 - total_gap) / n
        gap = 14
    card_h = 97
    y = BODY_START_Y + 15

    for i, m in enumerate(metrics[:n]):
        x = 29 + i * (card_w + gap)
        # Navy card
        add_round_rect(slide, x, y, card_w, card_h, NAVY)
        # Gold accent bar
        add_rect(slide, x, y, card_w, 3, GOLD)
        # Number — adaptive size
        val_str = str(m.get("value", ""))
        num_sz = 28 if len(val_str) <= 5 else 22
        add_textbox(slide, x, y + 15, card_w, 35, val_str,
                    font_size=num_sz, color=GOLD, font=FONT_HEADING, bold=True, align="ctr")
        # Label
        add_textbox(slide, x + 8, y + 55, card_w - 16, 18,
                    m.get("label", ""),
                    font_size=11, color=WHITE, font=FONT_BODY, align="ctr")
        # Sublabel
        if m.get("sublabel"):
            add_textbox(slide, x + 8, y + 75, card_w - 16, 16,
                        m["sublabel"],
                        font_size=9, color=GOLD, font=FONT_BODY, align="ctr")

    add_footer(slide, client_name, doc_label)
    return slide


def build_process_flow(prs, title, steps, client_name, doc_label):
    """Process flow — numbered circles with arrows and cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    n = min(len(steps), 5)
    gap = 11.3
    col_w = (CONTENT_W - 14 - gap * (n - 1)) / n
    circle_d = 40
    circle_y = BODY_START_Y + 8
    card_top = BODY_START_Y + 65
    card_h = BODY_END_Y - card_top - 5

    for i, step in enumerate(steps[:n]):
        col_x = 29 + i * (col_w + gap)
        circle_x = col_x + (col_w - circle_d) / 2

        # Numbered circle
        fill = GOLD if i == 0 else NAVY
        add_circle(slide, circle_x, circle_y, circle_d, fill,
                   str(i + 1), WHITE, 14, FONT_HEADING)

        # Arrow (except last)
        if i < n - 1:
            ax = circle_x + circle_d + 3
            next_cx = 29 + (i+1) * (col_w + gap) + (col_w - circle_d) / 2
            aw = next_cx - ax - 3
            ay = circle_y + circle_d / 2 - 4
            add_arrow(slide, ax, ay, aw, 8, GOLD)

        # Vertical connector
        add_line(slide, col_x + col_w / 2 - 0.5, circle_y + circle_d,
                 1, card_top - circle_y - circle_d, GOLD)

        # Card
        add_round_rect(slide, col_x, card_top, col_w, card_h, CARD_BG, GOLD)
        # Gold accent bar
        add_rect(slide, col_x, card_top, col_w, 3, GOLD)

        # Card text
        pad = CARD_PADDING
        add_multi_text(slide, col_x + pad, card_top + pad + 3,
                       col_w - 2*pad, card_h - 2*pad, [
            {'text': step.get('title', ''), 'size': 11, 'color': NAVY,
             'font': FONT_BODY, 'bold': True, 'spcAft': 500},
            {'text': step.get('body', step.get('description', '')),
             'size': 10, 'color': BODY_CLR, 'font': FONT_BODY, 'lnSpc': 1300},
        ])

    add_footer(slide, client_name, doc_label)
    return slide


def build_hub_spoke(prs, title, hub_text, cards, client_name, doc_label):
    """Hub-and-spoke — center circle with 4 surrounding cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    # Hub circle
    hub_r = 52
    hub_cx, hub_cy = 360, 222
    add_circle(slide, hub_cx - hub_r, hub_cy - hub_r, hub_r * 2, NAVY,
               hub_text, GOLD, 10, FONT_HEADING)

    # 4 cards: TL, TR, BL, BR
    positions = [
        (29, 75, 218, 130),    # TL
        (473, 75, 218, 130),   # TR
        (29, 235, 218, 130),   # BL
        (473, 235, 218, 130),  # BR
    ]

    for i, card in enumerate(cards[:4]):
        cx, cy, cw, ch = positions[i]
        add_round_rect(slide, cx, cy, cw, ch, CARD_BG, GOLD)
        add_rect(slide, cx, cy, cw, 3, GOLD)

        pad = CARD_PADDING
        paras = []
        if card.get('metric'):
            paras.append({'text': card['metric'], 'size': 10, 'color': GOLD,
                         'font': FONT_HEADING, 'bold': True, 'spcAft': 200})
        paras.append({'text': card.get('title', ''), 'size': 10, 'color': NAVY,
                     'font': FONT_BODY, 'bold': True, 'spcAft': 300})
        paras.append({'text': card.get('description', card.get('body', '')),
                     'size': 9, 'color': BODY_CLR, 'font': FONT_BODY, 'lnSpc': 1100})

        add_multi_text(slide, cx + pad, cy + pad + 3, cw - 2*pad, ch - 2*pad, paras)

    # Dashed connectors (simplified as thin lines with alpha)
    for cx, cy, cw, ch in positions:
        card_center_x = cx + cw / 2
        card_center_y = cy + ch / 2
        # Draw line from card edge toward hub
        add_line(slide, min(card_center_x, hub_cx), min(card_center_y, hub_cy),
                 abs(card_center_x - hub_cx), 1, GOLD, "50000")

    add_footer(slide, client_name, doc_label)
    return slide


def build_comparison(prs, title, left_data, right_data, client_name, doc_label):
    """Comparison matrix — two sides."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    mid = SLIDE_W / 2
    header_y = BODY_START_Y + 5
    header_h = 35
    content_y = header_y + header_h + 3
    content_h = BODY_END_Y - content_y - 5

    # Left header
    add_rect(slide, 29, header_y, mid - 35, header_h, NAVY)
    add_textbox(slide, 29, header_y + 8, mid - 35, 20,
                left_data.get("heading", "BEFORE"), font_size=12,
                color=WHITE, font=FONT_BODY, bold=True, align="ctr")

    # Right header
    add_rect(slide, mid + 5, header_y, mid - 35, header_h, GOLD)
    add_textbox(slide, mid + 5, header_y + 8, mid - 35, 20,
                right_data.get("heading", "AFTER"), font_size=12,
                color=WHITE, font=FONT_BODY, bold=True, align="ctr")

    # Left content
    add_round_rect(slide, 29, content_y, mid - 35, content_h, CARD_BG)
    paras_l = [{'text': p, 'size': 10, 'color': BODY_CLR, 'font': FONT_BODY,
                'lnSpc': 1400, 'spcAft': 400, 'bullet': '\u2022',
                'marL': '228600', 'indent': '-228600'}
               for p in left_data.get("paragraphs", [])]
    add_multi_text(slide, 39, content_y + 10, mid - 55, content_h - 20, paras_l)

    # Right content
    add_round_rect(slide, mid + 5, content_y, mid - 35, content_h, CARD_BG)
    paras_r = [{'text': p, 'size': 10, 'color': BODY_CLR, 'font': FONT_BODY,
                'lnSpc': 1400, 'spcAft': 400, 'bullet': '\u2022',
                'marL': '228600', 'indent': '-228600'}
               for p in right_data.get("paragraphs", [])]
    add_multi_text(slide, mid + 15, content_y + 10, mid - 55, content_h - 20, paras_r)

    # Divider
    add_line(slide, mid, header_y, 2, content_h + header_h + 3, NAVY, "50000")

    add_footer(slide, client_name, doc_label)
    return slide


def build_icon_cards(prs, title, items, client_name, doc_label, intro=None):
    """Grid of numbered icon cards (2 cols x N rows)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    y_start = BODY_START_Y + 8
    if intro:
        add_textbox(slide, MARGIN_L + 7, y_start, CONTENT_W - 14, 20, intro,
                    font_size=10, color=BODY_CLR, font=FONT_BODY)
        y_start += 25

    cols = 2
    gap_x, gap_y = 15, 7
    card_w = (CONTENT_W - 14 - gap_x * (cols - 1)) / cols
    card_h = 83
    circle_d = 25

    for i, item in enumerate(items[:6]):
        col = i % cols
        row = i // cols
        x = 25 + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)

        # Card with gold border
        add_round_rect(slide, x, y, card_w, card_h, WHITE, GOLD)

        # Numbered circle
        add_circle(slide, x + 7, y + 11, circle_d, GOLD,
                   str(i + 1), WHITE, 11, FONT_BODY)

        # Title right of circle
        add_textbox(slide, x + 40, y + 11, card_w - 50, 18,
                    item.get('title', ''), font_size=12, color=NAVY,
                    font=FONT_BODY, bold=True)

        # Description
        add_textbox(slide, x + 40, y + 33, card_w - 50, card_h - 40,
                    item.get('description', item.get('body', '')),
                    font_size=10, color=BODY_CLR, font=FONT_BODY)

    add_footer(slide, client_name, doc_label)
    return slide


def build_timeline(prs, title, milestones, client_name, doc_label):
    """Timeline with dots and cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    n = min(len(milestones), 5)
    timeline_y = BODY_START_Y + 20
    dot_sz = 16
    line_y = timeline_y + dot_sz / 2
    card_top = timeline_y + dot_sz + 25
    card_h = BODY_END_Y - card_top - 10

    # Horizontal gold line
    add_line(slide, 29, line_y, CONTENT_W - 14, 2, GOLD)

    spacing = (CONTENT_W - 14) / max(n - 1, 1) if n > 1 else 0
    card_w = (CONTENT_W - 14 - (n-1) * 8) / n if n > 1 else CONTENT_W - 14

    for i, ms in enumerate(milestones[:n]):
        dot_x = 29 + i * spacing - dot_sz / 2 if n > 1 else 29
        col_x = 29 + i * (card_w + 8) if n > 1 else 29

        # Dot
        fill = GOLD if i == 0 or i == n - 1 else NAVY
        add_circle(slide, dot_x, timeline_y, dot_sz, fill, "", WHITE, 8)

        # Year label above
        add_textbox(slide, dot_x - 15, timeline_y - 18, dot_sz + 30, 16,
                    str(ms.get('year', ms.get('label', ''))),
                    font_size=10, color=NAVY, font=FONT_BODY, bold=True, align="ctr")

        # Vertical connector
        add_line(slide, col_x + card_w / 2, timeline_y + dot_sz,
                 1, card_top - timeline_y - dot_sz, GOLD)

        # Card
        add_round_rect(slide, col_x, card_top, card_w, card_h, CARD_BG, GOLD)
        add_rect(slide, col_x, card_top, card_w, 3, GOLD)

        pad = CARD_PADDING
        add_multi_text(slide, col_x + pad, card_top + pad + 3,
                       card_w - 2*pad, card_h - 2*pad, [
            {'text': ms.get('title', ''), 'size': 10, 'color': NAVY,
             'font': FONT_BODY, 'bold': True, 'spcAft': 400},
            {'text': ms.get('description', ms.get('body', '')),
             'size': 9, 'color': BODY_CLR, 'font': FONT_BODY, 'lnSpc': 1200},
        ])

    add_footer(slide, client_name, doc_label)
    return slide


def build_quote(prs, quote_text, attribution, client_name, doc_label):
    """Quote slide — dark bg, gold quote mark."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Gold quote mark
    add_textbox(slide, MARGIN_L, 80, 60, 60, "\u201C",
                font_size=72, color=GOLD, font=FONT_HEADING)

    # Quote
    add_textbox(slide, 80, 140, CONTENT_W - 80, 140, quote_text,
                font_size=16, color=WHITE, font=FONT_HEADING, italic=True,
                line_spacing=2400)

    # Attribution
    if attribution:
        add_textbox(slide, 80, 300, CONTENT_W - 80, 25,
                    f"\u2014 {attribution}",
                    font_size=11, color=CREAM, font=FONT_BODY)

    add_footer(slide, client_name, doc_label)
    return slide


def build_table_slide(prs, title, headers, rows, client_name, doc_label):
    """Professional table slide using python-pptx table API with correct styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_y = BODY_START_Y + 8
    table_h = min(n_rows * 32 + 4, BODY_END_Y - table_y - 10)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Pt(29), Pt(table_y), Pt(CONTENT_W - 14), Pt(table_h)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = C_WHITE
                run.font.bold = True
                run.font.name = FONT_BODY

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF5, 0xF2)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
                    run.font.name = FONT_BODY
                    if j == 0:
                        run.font.bold = True

    add_footer(slide, client_name, doc_label)
    return slide


def build_closing(prs, client_name, doc_label, message="Thank you"):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_gold_rule(slide, (SLIDE_W - 288) / 2, 155, 288)
    add_textbox(slide, MARGIN_L, 170, CONTENT_W, 50, message,
                font_size=28, color=WHITE, font=FONT_HEADING, bold=True, align="ctr")
    add_textbox(slide, MARGIN_L, 230, CONTENT_W, 25,
                f"Prepared for {client_name}",
                font_size=14, color=CREAM, font=FONT_BODY, italic=True, align="ctr")
    add_textbox(slide, MARGIN_L, 260, CONTENT_W, 20,
                datetime.now().strftime("%B %Y"),
                font_size=10, color=CREAM, font=FONT_BODY, align="ctr")
    add_gold_rule(slide, (SLIDE_W - 288) / 2, 290, 288)

    add_footer(slide, client_name, doc_label)
    return slide


# ============================================================
# ASSEMBLER
# ============================================================
def assemble(content, doc_type="methodology"):
    prs = Presentation()
    prs.slide_width = Pt(SLIDE_W)
    prs.slide_height = Pt(SLIDE_H)

    c = content
    client = c.get("client_name", "Client")
    visa = c.get("visa_type", "")
    label = c.get("doc_label", "Professional Dossier")
    title = c.get("title", "Document")
    subtitle = c.get("subtitle", "")
    sections = c.get("sections", [])

    # Cover
    build_cover(prs, client, title, subtitle, label, visa)

    # TOC
    if sections:
        build_toc(prs, sections, client, label)

    # Sections
    for si, section in enumerate(sections):
        build_section_divider(prs, section.get("title", ""),
                             section.get("subtitle", ""), client, label,
                             number=str(si + 1))

        for slide_data in section.get("slides", []):
            t = slide_data.get("type", "content")

            if t == "content":
                build_content(prs, slide_data.get("title", ""),
                             slide_data.get("paragraphs", []),
                             client, label,
                             bullets=slide_data.get("bullets"))

            elif t == "metrics":
                build_metrics(prs, slide_data.get("title", ""),
                             slide_data.get("metrics", []),
                             client, label)

            elif t == "process_flow":
                build_process_flow(prs, slide_data.get("title", ""),
                                  slide_data.get("steps", []),
                                  client, label)

            elif t == "hub_spoke":
                build_hub_spoke(prs, slide_data.get("title", ""),
                               slide_data.get("hub_text", "METHOD"),
                               slide_data.get("cards", []),
                               client, label)

            elif t == "comparison" or t == "two_column":
                build_comparison(prs, slide_data.get("title", ""),
                                slide_data.get("left", {}),
                                slide_data.get("right", {}),
                                client, label)

            elif t == "icon_cards" or t == "icon_grid" or t == "icon_list":
                build_icon_cards(prs, slide_data.get("title", ""),
                                slide_data.get("items", []),
                                client, label,
                                intro=slide_data.get("intro_text"))

            elif t == "timeline":
                build_timeline(prs, slide_data.get("title", ""),
                              slide_data.get("milestones", []),
                              client, label)

            elif t == "table":
                build_table_slide(prs, slide_data.get("title", ""),
                                 slide_data.get("headers", []),
                                 slide_data.get("rows", []),
                                 client, label)

            elif t == "quote":
                build_quote(prs, slide_data.get("quote", ""),
                           slide_data.get("attribution", ""),
                           client, label)

    # Closing
    build_closing(prs, client, label, c.get("closing_message", "Thank you"))

    return prs


# ============================================================
# CLI
# ============================================================
def main():
    parser = argparse.ArgumentParser(description="Petition Engine PPTX V2")
    parser.add_argument("--content", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--type", default="methodology")
    args = parser.parse_args()

    with open(args.content, 'r', encoding='utf-8') as f:
        content = json.load(f)

    print(f"Generating {args.type} PPTX V2...")
    prs = assemble(content, args.type)

    os.makedirs(os.path.dirname(args.output) or '.', exist_ok=True)
    prs.save(args.output)
    print(f"DONE: {args.output}")
    print(f"  Slides: {len(prs.slides)} | Size: {os.path.getsize(args.output)/1024:.0f} KB")


if __name__ == "__main__":
    main()
