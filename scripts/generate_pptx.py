#!/usr/bin/env python3
"""
PETITION ENGINE — Professional PPTX Generator
Design DNA: André Cerbasi Dossier (39-slide benchmark)
Palette: Navy #1B2A4A | Gold #B8860B | Cream #E8D5B7 | Dark Gray #3A3A3A

Usage:
    python3 generate_pptx.py --type methodology --content content.json --output output.pptx
    python3 generate_pptx.py --type declaration --content content.json --output output.pptx
"""

import json
import os
import sys
import argparse
import re
import glob as globmod
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from datetime import datetime

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ICONS_DIR = os.path.join(SCRIPT_DIR, 'icons')

# ============================================================
# DESIGN SYSTEM — Cerbasi DNA
# ============================================================
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
GOLD      = RGBColor(0xB8, 0x86, 0x0B)
CREAM     = RGBColor(0xE8, 0xD5, 0xB7)
DARK_GRAY = RGBColor(0x3A, 0x3A, 0x3A)
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MED_GRAY  = RGBColor(0x99, 0x99, 0x99)
DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)

# Typography
FONT_TITLE    = "Georgia"
FONT_BODY     = "Calibri"
FONT_ACCENT   = "Georgia"

# Slide dimensions (widescreen 10x5.625)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Margins
M_LEFT   = Inches(0.6)
M_RIGHT  = Inches(0.6)
M_TOP    = Inches(0.5)
M_BOTTOM = Inches(0.5)

CONTENT_W = SLIDE_W - M_LEFT - M_RIGHT  # usable width


# ============================================================
# HELPER FUNCTIONS
# ============================================================
def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, opacity=1.0):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_text(tf, text, size=Pt(11), color=DARK_GRAY, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name=None, space_after=Pt(4)):
    """Add a paragraph with formatted text to a text_frame."""
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    if tf.paragraphs[0].text:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name or FONT_BODY
    return p


def add_footer(slide, client_name, doc_label):
    """Add footer bar to bottom of slide — cream text on dark bar."""
    # Dark bar
    bar_h = Inches(0.3)
    bar_top = SLIDE_H - bar_h
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_W, bar_h, NAVY)

    # Footer text
    tf = add_textbox(slide, M_LEFT, bar_top, CONTENT_W, bar_h)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"{client_name} | {doc_label}"
    run.font.size = Pt(8)
    run.font.color.rgb = CREAM
    run.font.name = FONT_BODY
    run.font.italic = True


def add_gold_accent_line(slide, top, width=None):
    """Add a thin gold horizontal line."""
    w = width or Inches(2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M_LEFT, top, w, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE BUILDERS
# ============================================================
def build_cover_slide(prs, client_name, title, subtitle, doc_label, visa_type=""):
    """Slide 1 — Cover with dark background, gold accents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Gold accent bar top
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), GOLD)

    # Client name — large bold
    tf = add_textbox(slide, M_LEFT, Inches(1.2), CONTENT_W, Inches(0.8))
    add_text(tf, client_name, size=Pt(36), color=WHITE, bold=True,
             font_name=FONT_TITLE, alignment=PP_ALIGN.LEFT)

    # Visa type — small
    if visa_type:
        tf2 = add_textbox(slide, M_LEFT, Inches(1.9), CONTENT_W, Inches(0.4))
        add_text(tf2, visa_type, size=Pt(16), color=GOLD, bold=False,
                 font_name=FONT_BODY, alignment=PP_ALIGN.LEFT)

    # Gold line
    add_gold_accent_line(slide, Inches(2.3), Inches(3))

    # Title
    tf3 = add_textbox(slide, M_LEFT, Inches(2.6), CONTENT_W, Inches(0.6))
    add_text(tf3, title, size=Pt(18), color=CREAM, bold=False,
             font_name=FONT_ACCENT, alignment=PP_ALIGN.LEFT, italic=True)

    # Subtitle
    tf4 = add_textbox(slide, M_LEFT, Inches(3.2), CONTENT_W, Inches(0.6))
    add_text(tf4, subtitle, size=Pt(11), color=LIGHT_GRAY, bold=False,
             font_name=FONT_BODY, alignment=PP_ALIGN.LEFT)

    # Date
    tf5 = add_textbox(slide, M_LEFT, Inches(4.4), CONTENT_W, Inches(0.3))
    add_text(tf5, datetime.now().strftime("%B %Y"), size=Pt(10), color=MED_GRAY,
             font_name=FONT_BODY, alignment=PP_ALIGN.LEFT)

    add_footer(slide, client_name, doc_label)
    return slide


def build_toc_slide(prs, sections, client_name, doc_label):
    """Slide 2 — Table of Contents with Roman numerals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.6))
    add_text(tf_title, "Table of Contents", size=Pt(32), color=NAVY, bold=True,
             font_name=FONT_TITLE)

    # Gold line under title
    add_gold_accent_line(slide, Inches(1.0), Inches(2))

    # Entries — 2 columns if > 5
    romans = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
              "XI", "XII", "XIII", "XIV", "XV"]

    col_w = CONTENT_W / 2 if len(sections) > 5 else CONTENT_W
    y_start = Inches(1.3)
    items_per_col = (len(sections) + 1) // 2 if len(sections) > 5 else len(sections)

    for i, sec in enumerate(sections):
        col = 0 if i < items_per_col else 1
        row = i if col == 0 else i - items_per_col
        x = M_LEFT + (col * col_w)
        y = y_start + Inches(row * 0.45)

        # Roman numeral + section title in ONE textbox to prevent overlap
        tf_entry = add_textbox(slide, x, y, col_w - Inches(0.1), Inches(0.4))
        p = tf_entry.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # Roman numeral run
        run_num = p.add_run()
        run_num.text = (romans[i] if i < len(romans) else str(i+1)) + "   "
        run_num.font.size = Pt(16)
        run_num.font.color.rgb = GOLD
        run_num.font.bold = True
        run_num.font.name = FONT_TITLE
        # Section title run
        run_title = p.add_run()
        run_title.text = sec.get("title", f"Section {i+1}")
        run_title.font.size = Pt(12)
        run_title.font.color.rgb = NAVY
        run_title.font.bold = True
        run_title.font.name = FONT_BODY

    add_footer(slide, client_name, doc_label)
    return slide


def build_section_divider(prs, section_title, section_subtitle, client_name, doc_label):
    """Full-bleed section divider — gold title on dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Gold accent bar
    add_shape_rect(slide, Inches(0.6), Inches(2.0), Inches(3), Pt(3), GOLD)

    # Title — large gold
    tf = add_textbox(slide, M_LEFT, Inches(2.2), CONTENT_W, Inches(1.0))
    add_text(tf, section_title.upper(), size=Pt(44), color=GOLD, bold=True,
             font_name=FONT_TITLE, alignment=PP_ALIGN.LEFT)

    # Subtitle
    if section_subtitle:
        tf2 = add_textbox(slide, M_LEFT, Inches(3.3), CONTENT_W, Inches(0.5))
        add_text(tf2, section_subtitle, size=Pt(14), color=CREAM, italic=True,
                 font_name=FONT_BODY, alignment=PP_ALIGN.LEFT)

    add_footer(slide, client_name, doc_label)
    return slide


def build_content_slide(prs, title, body_paragraphs, client_name, doc_label,
                        bullets=None, highlight_box=None):
    """Standard content slide — white bg, navy title, body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title — positioned with room below
    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.55))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True,
             font_name=FONT_TITLE)

    # Gold accent under title — clear of text
    add_gold_accent_line(slide, Inches(1.05), Inches(1.5))

    # Body text — starts below the gold line
    y_body = Inches(1.2)
    body_w = CONTENT_W
    if highlight_box:
        body_w = CONTENT_W * 0.6  # leave room for highlight

    tf_body = add_textbox(slide, M_LEFT, y_body, body_w, Inches(3.5))

    for para_text in body_paragraphs:
        if not para_text.strip():
            continue
        add_text(tf_body, para_text.strip(), size=Pt(11), color=DARK_GRAY,
                 font_name=FONT_BODY, space_after=Pt(8))

    # Bullet points
    if bullets:
        for bullet in bullets:
            p = tf_body.add_paragraph()
            p.space_after = Pt(4)
            p.level = 0
            run = p.add_run()
            run.text = f"  {bullet.strip()}"
            run.font.size = Pt(11)
            run.font.color.rgb = DARK_GRAY
            run.font.name = FONT_BODY

    # Highlight box (right side)
    if highlight_box:
        box_left = M_LEFT + body_w + Inches(0.3)
        box_w = CONTENT_W - body_w - Inches(0.3)
        rect = add_shape_rect(slide, box_left, y_body, box_w, Inches(3.0), LIGHT_GRAY)

        tf_box = add_textbox(slide, box_left + Inches(0.15), y_body + Inches(0.15),
                             box_w - Inches(0.3), Inches(2.7))
        if highlight_box.get("title"):
            add_text(tf_box, highlight_box["title"], size=Pt(12), color=NAVY,
                     bold=True, font_name=FONT_BODY)
        if highlight_box.get("text"):
            add_text(tf_box, highlight_box["text"], size=Pt(10), color=DARK_GRAY,
                     font_name=FONT_BODY)

    add_footer(slide, client_name, doc_label)
    return slide


def build_key_metrics_slide(prs, title, metrics, client_name, doc_label):
    """Metrics slide — big numbers with labels. Each metric in its own card, no overlap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True, font_name=FONT_TITLE)
    add_gold_accent_line(slide, Inches(1.0), Inches(1.5))

    num = min(len(metrics), 4)
    col_w = CONTENT_W / num
    gap = Inches(0.15)
    card_w = col_w - gap * 2
    y_card = Inches(1.4)
    card_h = Inches(2.5)

    for i, metric in enumerate(metrics[:4]):
        x = M_LEFT + (i * col_w) + gap
        # Card background
        add_shape_rect(slide, x, y_card, card_w, card_h, LIGHT_GRAY)
        # Big number — centered in card
        tf_num = add_textbox(slide, x, y_card + Inches(0.4), card_w, Inches(0.8))
        add_text(tf_num, str(metric.get("value", "N/A")), size=Pt(40), color=GOLD,
                 bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)
        # Label — below number
        tf_label = add_textbox(slide, x + Inches(0.1), y_card + Inches(1.4),
                               card_w - Inches(0.2), Inches(0.9))
        add_text(tf_label, metric.get("label", ""), size=Pt(10), color=NAVY,
                 bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, client_name, doc_label)
    return slide


def build_table_slide(prs, title, headers, rows, client_name, doc_label):
    """Table slide — navy header, clean rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(20), color=NAVY, bold=True,
             font_name=FONT_TITLE)

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    table_top = Inches(1.0)
    table_h = Inches(min(0.35 * num_rows, 4.0))

    table_shape = slide.shapes.add_table(num_rows, num_cols,
                                          M_LEFT, table_top, CONTENT_W, table_h)
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.name = FONT_BODY

    # Data rows
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            if j >= num_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            # Remove lateral borders via XML
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = DARK_GRAY
                    run.font.name = FONT_BODY
                    if j == 0:
                        run.font.bold = True

    add_footer(slide, client_name, doc_label)
    return slide


def build_two_column_slide(prs, title, left_content, right_content, client_name, doc_label):
    """Two-column layout for comparisons, before/after, dual concepts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(24), color=NAVY, bold=True,
             font_name=FONT_TITLE)

    add_gold_accent_line(slide, Inches(0.9), Inches(1.5))

    col_w = (CONTENT_W - Inches(0.3)) / 2
    y_cols = Inches(1.2)

    # Left column
    tf_left = add_textbox(slide, M_LEFT, y_cols, col_w, Inches(3.5))
    if left_content.get("heading"):
        add_text(tf_left, left_content["heading"], size=Pt(16), color=NAVY,
                 bold=True, font_name=FONT_BODY)
    for para in left_content.get("paragraphs", []):
        add_text(tf_left, para, size=Pt(11), color=DARK_GRAY, font_name=FONT_BODY,
                 space_after=Pt(6))

    # Right column
    tf_right = add_textbox(slide, M_LEFT + col_w + Inches(0.3), y_cols, col_w, Inches(3.5))
    if right_content.get("heading"):
        add_text(tf_right, right_content["heading"], size=Pt(16), color=NAVY,
                 bold=True, font_name=FONT_BODY)
    for para in right_content.get("paragraphs", []):
        add_text(tf_right, para, size=Pt(11), color=DARK_GRAY, font_name=FONT_BODY,
                 space_after=Pt(6))

    add_footer(slide, client_name, doc_label)
    return slide


def build_quote_slide(prs, quote_text, attribution, client_name, doc_label):
    """Full-slide quote — elegant, centered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Quote mark
    tf_mark = add_textbox(slide, M_LEFT, Inches(1.0), Inches(1), Inches(1))
    add_text(tf_mark, "\u201C", size=Pt(72), color=GOLD, font_name=FONT_TITLE,
             alignment=PP_ALIGN.LEFT)

    # Quote text
    tf_quote = add_textbox(slide, Inches(1.2), Inches(1.8), Inches(7.5), Inches(2.0))
    add_text(tf_quote, quote_text, size=Pt(18), color=WHITE, italic=True,
             font_name=FONT_ACCENT, alignment=PP_ALIGN.LEFT, space_after=Pt(12))

    # Attribution
    if attribution:
        tf_attr = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(7.5), Inches(0.4))
        add_text(tf_attr, f"-- {attribution}", size=Pt(12), color=CREAM,
                 font_name=FONT_BODY, alignment=PP_ALIGN.LEFT)

    add_footer(slide, client_name, doc_label)
    return slide


def build_closing_slide(prs, client_name, doc_label, message="Thank you"):
    """Final slide — dark bg, centered message."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Gold line
    line_w = Inches(3)
    add_shape_rect(slide, (SLIDE_W - line_w) / 2, Inches(2.0), line_w, Pt(3), GOLD)

    # Message
    tf = add_textbox(slide, M_LEFT, Inches(2.3), CONTENT_W, Inches(1.0))
    add_text(tf, message, size=Pt(32), color=WHITE, bold=True,
             font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

    # Subtext
    tf2 = add_textbox(slide, M_LEFT, Inches(3.3), CONTENT_W, Inches(0.5))
    add_text(tf2, f"Prepared for {client_name}", size=Pt(14), color=CREAM,
             font_name=FONT_BODY, alignment=PP_ALIGN.CENTER, italic=True)

    # Date
    tf3 = add_textbox(slide, M_LEFT, Inches(3.8), CONTENT_W, Inches(0.3))
    add_text(tf3, datetime.now().strftime("%B %Y"), size=Pt(10), color=MED_GRAY,
             font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)

    # Gold line bottom
    add_shape_rect(slide, (SLIDE_W - line_w) / 2, Inches(4.3), line_w, Pt(3), GOLD)

    add_footer(slide, client_name, doc_label)
    return slide


# ============================================================
# PHOTO & ICON HELPERS
# ============================================================
def find_client_photo(client_docs_path):
    """Search client folder for a profile photo. Returns path or None."""
    if not client_docs_path or not os.path.exists(client_docs_path):
        return None
    # Common photo patterns
    patterns = ['**/foto*', '**/photo*', '**/profile*', '**/headshot*',
                '**/retrato*', '**/*perfil*', '**/*passport*']
    extensions = ['.jpg', '.jpeg', '.png', '.webp']
    for pat in patterns:
        for ext in extensions:
            matches = globmod.glob(os.path.join(client_docs_path, pat + ext), recursive=True)
            if matches:
                return matches[0]
    # Fallback: any image in root
    for ext in extensions:
        matches = globmod.glob(os.path.join(client_docs_path, '*' + ext))
        if matches:
            return matches[0]
    return None

def get_icon_path(icon_name):
    """Get path to an icon PNG. Falls back to methodology icon."""
    path = os.path.join(ICONS_DIR, f"{icon_name}.png")
    if os.path.exists(path):
        return path
    # Try to generate it
    try:
        from pptx_icons import get_icon
        return get_icon(icon_name)
    except:
        pass
    # Fallback
    fallback = os.path.join(ICONS_DIR, "methodology.png")
    return fallback if os.path.exists(fallback) else None


# ============================================================
# VISUAL SLIDE BUILDERS (Gamma-level)
# ============================================================
def build_process_flow_slide(prs, title, steps, client_name, doc_label):
    """Process flow with chevron arrows — strict column isolation to prevent overlap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True, font_name=FONT_TITLE)
    add_gold_accent_line(slide, Inches(1.0), Inches(1.5))

    num_steps = min(len(steps), 5)
    gap = Inches(0.1)
    step_w = (CONTENT_W - gap * (num_steps - 1)) / num_steps
    y_icon = Inches(1.3)
    icon_sz = Inches(0.45)
    y_chevron = Inches(1.9)
    chevron_h = Inches(0.45)
    y_title = Inches(2.5)
    y_desc = Inches(2.9)

    for i, step in enumerate(steps[:5]):
        x = M_LEFT + i * (step_w + gap)

        # Icon centered above chevron
        icon_name = step.get('icon', 'process')
        icon_path = get_icon_path(icon_name)
        if icon_path:
            try:
                ix = x + (step_w - icon_sz) / 2
                slide.shapes.add_picture(icon_path, int(ix), int(y_icon), icon_sz, icon_sz)
            except:
                pass

        # Chevron shape — exact width, no bleed
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, int(x), int(y_chevron), int(step_w), int(chevron_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()

        # Step title — STRICTLY within column width
        tf_step = add_textbox(slide, x, y_title, step_w, Inches(0.35))
        add_text(tf_step, step.get('title', ''), size=Pt(10), color=NAVY,
                 bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)

        # Description — STRICTLY within column width
        tf_desc = add_textbox(slide, x, y_desc, step_w, Inches(1.8))
        add_text(tf_desc, step.get('description', ''), size=Pt(9), color=DARK_GRAY,
                 font_name=FONT_BODY, alignment=PP_ALIGN.CENTER, space_after=Pt(3))

    add_footer(slide, client_name, doc_label)
    return slide


def build_icon_grid_slide(prs, title, items, client_name, doc_label, intro_text=None):
    """3-column grid with icons — strict column isolation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True, font_name=FONT_TITLE)

    y_start = Inches(1.1)
    if intro_text:
        tf_intro = add_textbox(slide, M_LEFT, y_start, CONTENT_W, Inches(0.45))
        add_text(tf_intro, intro_text, size=Pt(10), color=DARK_GRAY, font_name=FONT_BODY)
        y_start += Inches(0.5)

    cols = min(len(items), 3)
    gap = Inches(0.15)
    col_w = (CONTENT_W - gap * (cols - 1)) / cols
    icon_sz = Inches(0.4)
    card_h = Inches(1.5)

    for i, item in enumerate(items[:6]):
        col = i % cols
        row = i // cols
        x = M_LEFT + col * (col_w + gap)
        y = y_start + row * (card_h + Inches(0.15))

        # Card background — exact width, no bleed
        add_shape_rect(slide, x, y, col_w, card_h, LIGHT_GRAY)

        # Icon — top-left of card
        icon_name = item.get('icon', 'methodology')
        icon_path = get_icon_path(icon_name)
        if icon_path:
            try:
                slide.shapes.add_picture(icon_path, int(x + Inches(0.12)),
                                         int(y + Inches(0.12)), icon_sz, icon_sz)
            except:
                pass

        # Title — within card bounds
        tf_item = add_textbox(slide, x + Inches(0.1), y + Inches(0.6),
                             col_w - Inches(0.2), Inches(0.25))
        add_text(tf_item, item.get('title', ''), size=Pt(10), color=NAVY,
                 bold=True, font_name=FONT_BODY)

        # Description — within card bounds
        tf_desc = add_textbox(slide, x + Inches(0.1), y + Inches(0.85),
                             col_w - Inches(0.2), Inches(0.55))
        add_text(tf_desc, item.get('description', ''), size=Pt(8.5),
                 color=DARK_GRAY, font_name=FONT_BODY, space_after=Pt(2))

    add_footer(slide, client_name, doc_label)
    return slide


def build_photo_content_slide(prs, title, paragraphs, photo_path, client_name, doc_label,
                               photo_side='left'):
    """Split slide: photo on one side, text on the other — like Dabus formation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    photo_w = Inches(4.0)
    text_w = CONTENT_W - photo_w - Inches(0.3)

    if photo_side == 'left':
        photo_x = M_LEFT
        text_x = M_LEFT + photo_w + Inches(0.3)
    else:
        text_x = M_LEFT
        photo_x = M_LEFT + text_w + Inches(0.3)

    # Photo
    if photo_path and os.path.exists(photo_path):
        try:
            slide.shapes.add_picture(photo_path, int(photo_x), int(M_TOP),
                                     photo_w, Inches(4.5))
        except Exception as e:
            # Fallback: gray placeholder
            add_shape_rect(slide, photo_x, M_TOP, photo_w, Inches(4.5), LIGHT_GRAY)
    else:
        add_shape_rect(slide, photo_x, M_TOP, photo_w, Inches(4.5), LIGHT_GRAY)

    # Title
    tf_title = add_textbox(slide, text_x, M_TOP, text_w, Inches(0.7))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True, font_name=FONT_TITLE)

    # Body text
    tf_body = add_textbox(slide, text_x, Inches(1.3), text_w, Inches(3.5))
    for para in paragraphs[:4]:  # max 4 paragraphs
        add_text(tf_body, para.strip(), size=Pt(10), color=DARK_GRAY,
                 font_name=FONT_BODY, space_after=Pt(8))

    add_footer(slide, client_name, doc_label)
    return slide


def build_icon_list_slide(prs, title, items, client_name, doc_label, intro_text=None):
    """Vertical list with icons on left — like Dabus legal training / digital humanization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    tf_title = add_textbox(slide, M_LEFT, M_TOP, CONTENT_W, Inches(0.5))
    add_text(tf_title, title, size=Pt(22), color=NAVY, bold=True, font_name=FONT_TITLE)

    y_start = Inches(1.1)

    if intro_text:
        tf_intro = add_textbox(slide, M_LEFT, y_start, CONTENT_W, Inches(0.4))
        add_text(tf_intro, intro_text, size=Pt(10), color=DARK_GRAY, font_name=FONT_BODY)
        y_start += Inches(0.5)

    icon_size = Inches(0.4)
    item_h = Inches(0.85)

    for i, item in enumerate(items[:5]):  # max 5 items
        y = y_start + (i * item_h)

        # Icon
        icon_name = item.get('icon', 'process')
        icon_path = get_icon_path(icon_name)
        if icon_path:
            try:
                slide.shapes.add_picture(
                    icon_path, int(M_LEFT), int(y), icon_size, icon_size
                )
            except:
                pass

        # Title
        tf_item = add_textbox(slide, M_LEFT + Inches(0.6), y,
                             CONTENT_W - Inches(0.7), Inches(0.3))
        add_text(tf_item, item.get('title', ''), size=Pt(12), color=NAVY,
                 bold=True, font_name=FONT_BODY)

        # Description
        tf_desc = add_textbox(slide, M_LEFT + Inches(0.6), y + Inches(0.3),
                             CONTENT_W - Inches(0.7), Inches(0.45))
        add_text(tf_desc, item.get('description', ''), size=Pt(9),
                 color=DARK_GRAY, font_name=FONT_BODY)

    add_footer(slide, client_name, doc_label)
    return slide


# ============================================================
# DOCUMENT ASSEMBLER
# ============================================================
def assemble_presentation(content, doc_type):
    """
    Assemble a full presentation from structured content JSON.

    Expected content format:
    {
        "client_name": "Full Name",
        "visa_type": "EB-1A Extraordinary Ability",
        "doc_label": "Professional Methodology Dossier",
        "title": "Methodology — Comprehensive Analysis",
        "subtitle": "Detailed documentation of ...",
        "sections": [
            {
                "title": "Section Title",
                "subtitle": "Optional subtitle for divider",
                "slides": [
                    {
                        "type": "content|metrics|table|two_column|quote",
                        "title": "Slide Title",
                        "paragraphs": ["..."],
                        "bullets": ["..."],
                        "highlight_box": {"title": "...", "text": "..."},
                        "metrics": [{"value": "500+", "label": "..."}],
                        "headers": ["..."], "rows": [["..."]],
                        "left": {"heading": "...", "paragraphs": ["..."]},
                        "right": {"heading": "...", "paragraphs": ["..."]},
                        "quote": "...", "attribution": "..."
                    }
                ]
            }
        ],
        "closing_message": "Optional closing text"
    }
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    client_name = content.get("client_name", "Client")
    visa_type = content.get("visa_type", "")
    doc_label = content.get("doc_label", "Professional Dossier")
    title = content.get("title", "Document")
    subtitle = content.get("subtitle", "")
    sections = content.get("sections", [])
    client_docs_path = content.get("client_docs_path", "")

    # Find client photo
    client_photo = content.get("photo_path") or find_client_photo(client_docs_path)
    if client_photo:
        print(f"  Client photo found: {client_photo}")

    # 1. Cover slide
    build_cover_slide(prs, client_name, title, subtitle, doc_label, visa_type)

    # 2. TOC slide
    if sections:
        build_toc_slide(prs, sections, client_name, doc_label)

    # 3. Section slides
    for section in sections:
        # Section divider
        build_section_divider(prs, section.get("title", ""),
                             section.get("subtitle", ""), client_name, doc_label)

        # Section content slides
        for slide_data in section.get("slides", []):
            slide_type = slide_data.get("type", "content")

            if slide_type == "content":
                build_content_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("paragraphs", []),
                    client_name, doc_label,
                    bullets=slide_data.get("bullets"),
                    highlight_box=slide_data.get("highlight_box")
                )

            elif slide_type == "metrics":
                build_key_metrics_slide(
                    prs,
                    slide_data.get("title", "Key Metrics"),
                    slide_data.get("metrics", []),
                    client_name, doc_label
                )

            elif slide_type == "table":
                build_table_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("headers", []),
                    slide_data.get("rows", []),
                    client_name, doc_label
                )

            elif slide_type == "two_column":
                build_two_column_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("left", {}),
                    slide_data.get("right", {}),
                    client_name, doc_label
                )

            elif slide_type == "quote":
                build_quote_slide(
                    prs,
                    slide_data.get("quote", ""),
                    slide_data.get("attribution", ""),
                    client_name, doc_label
                )

            elif slide_type == "process_flow":
                build_process_flow_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("steps", []),
                    client_name, doc_label
                )

            elif slide_type == "icon_grid":
                build_icon_grid_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("items", []),
                    client_name, doc_label,
                    intro_text=slide_data.get("intro_text")
                )

            elif slide_type == "icon_list":
                build_icon_list_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("items", []),
                    client_name, doc_label,
                    intro_text=slide_data.get("intro_text")
                )

            elif slide_type == "photo_content":
                photo_path = slide_data.get("photo_path") or client_photo
                build_photo_content_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("paragraphs", []),
                    photo_path,
                    client_name, doc_label,
                    photo_side=slide_data.get("photo_side", "left")
                )

    # 4. Closing slide
    closing_msg = content.get("closing_message", "Thank you")
    build_closing_slide(prs, client_name, doc_label, closing_msg)

    return prs


# ============================================================
# CLI
# ============================================================
def main():
    parser = argparse.ArgumentParser(description="Petition Engine PPTX Generator")
    parser.add_argument("--content", required=True, help="Path to content JSON")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--type", default="methodology",
                        choices=["methodology", "declaration"],
                        help="Document type")
    args = parser.parse_args()

    if not os.path.exists(args.content):
        print(f"ERROR: Content file not found: {args.content}")
        sys.exit(1)

    with open(args.content, 'r', encoding='utf-8') as f:
        content = json.load(f)

    print(f"Generating {args.type} PPTX...")
    prs = assemble_presentation(content, args.type)

    os.makedirs(os.path.dirname(args.output) or '.', exist_ok=True)
    prs.save(args.output)

    slide_count = len(prs.slides)
    file_size = os.path.getsize(args.output)
    print(f"DONE: {args.output}")
    print(f"  Slides: {slide_count}")
    print(f"  Size: {file_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
